"""Extract Magna-Tiles API support docs into src/constants/knowledgeDocs.json."""

from __future__ import annotations

import json
import re
from pathlib import Path

from PyPDF2 import PdfReader
from docx import Document
from openpyxl import load_workbook
from pptx import Presentation

ROOT = Path(__file__).resolve().parents[1]
OUT = ROOT / "src" / "constants" / "knowledgeDocs.json"
MAGNA = Path(r"C:\Users\CihanHartamaci\Desktop\Projects\Magna-Tiles")
EXTRACT = ROOT / "_kb_extract"


def clean_text(value: str) -> str:
    text = value.replace("\x00", " ")
    previous = None
    while previous != text:
        previous = text
        text = re.sub(r"(\S)\n[ \t]*\n[ \t]*(\S)", r"\1 \2", text)
    text = re.sub(r"[ \t]+", " ", text)
    text = re.sub(r"\n{3,}", "\n\n", text)
    return text.strip()


def pdf_text(path: Path) -> str:
    reader = PdfReader(str(path))
    pages = []
    for index, page in enumerate(reader.pages, start=1):
        extracted = page.extract_text() or ""
        if extracted.strip():
            pages.append(f"--- Page {index} ---\n{extracted}")
    return clean_text("\n\n".join(pages))


def docx_text(path: Path) -> str:
    document = Document(str(path))
    parts = [paragraph.text for paragraph in document.paragraphs if paragraph.text.strip()]
    for table in document.tables:
        for row in table.rows:
            cells = [cell.text.strip() for cell in row.cells if cell.text.strip()]
            if cells:
                parts.append(" | ".join(cells))
    return clean_text("\n".join(parts))


def pptx_text(path: Path) -> str:
    presentation = Presentation(str(path))
    parts = []
    for index, slide in enumerate(presentation.slides, start=1):
        slide_bits = []
        for shape in slide.shapes:
            if getattr(shape, "has_text_frame", False) and shape.text_frame:
                text = shape.text_frame.text.strip()
                if text:
                    slide_bits.append(text)
            if getattr(shape, "has_table", False):
                table = shape.table
                for row in table.rows:
                    cells = [cell.text.strip() for cell in row.cells if cell.text.strip()]
                    if cells:
                        slide_bits.append(" | ".join(cells))
        if slide_bits:
            parts.append(f"--- Slide {index} ---\n" + "\n".join(slide_bits))
    return clean_text("\n\n".join(parts))


def xlsx_text(path: Path) -> str:
    workbook = load_workbook(str(path), data_only=True, read_only=True)
    parts = []
    for sheet in workbook.worksheets:
        rows = []
        for row in sheet.iter_rows(values_only=True):
            cells = [str(cell).strip() for cell in row if cell is not None and str(cell).strip()]
            if cells:
                rows.append(" | ".join(cells))
        if rows:
            parts.append(f"--- Sheet: {sheet.title} ---\n" + "\n".join(rows))
    workbook.close()
    return clean_text("\n\n".join(parts))


def json_text(path: Path) -> str:
    data = json.loads(path.read_text(encoding="utf-8"))
    pretty = json.dumps(data, indent=2, ensure_ascii=False)
    return clean_text(f"Example JSON payload from {path.name}:\n{pretty}")


EXTRACTORS = {
    ".pdf": pdf_text,
    ".docx": docx_text,
    ".pptx": pptx_text,
    ".xlsx": xlsx_text,
    ".json": json_text,
}


def article_from(path: Path, origin: str, url: str) -> dict | None:
    extractor = EXTRACTORS.get(path.suffix.lower())
    if not extractor:
        return None
    print(f"Extracting {path}")
    content = extractor(path)
    if len(content) < 80:
        print(f"  skipped (too short: {len(content)} chars)")
        return None
    title = path.stem.replace("_", " ").replace("-", " ").strip()
    title = re.sub(r"\s+", " ", title)
    if path.suffix.lower() == ".json":
        title = f"Example JSON: {title}"
    print(f"  {len(content)} chars")
    return {
        "title": title,
        "origin": origin,
        "filename": path.name,
        "url": url,
        "content": content,
    }


def main() -> None:
    jobs = []

    support_dir = EXTRACT / "API_Support_Doc"
    for path in sorted(support_dir.iterdir()):
        jobs.append(
            (
                path,
                "Magna-Tiles / API_Support_Doc.zip",
                f"kb://magna-tiles/API_Support_Doc/{path.name}",
            )
        )

    inventory_dir = EXTRACT / "Inventory_List_APIs_Guide"
    for path in sorted(inventory_dir.iterdir()):
        # Standalone webhook PDF is ingested separately to avoid duplicate chunks.
        if path.name.lower() == "shipment_inventory_webhook_guide.pdf":
            continue
        jobs.append(
            (
                path,
                "Magna-Tiles / Inventory_List_APIs_Guide.zip",
                f"kb://magna-tiles/Inventory_List_APIs_Guide/{path.name}",
            )
        )

    jobs.extend(
        [
            (
                MAGNA / "Shipment Order - Retailer Usage Documentation.pdf",
                "Magna-Tiles",
                "kb://magna-tiles/Shipment Order - Retailer Usage Documentation.pdf",
            ),
            (
                MAGNA / "Logiwa_IO_API_Carrier_Shipping_Option_Guide.docx",
                "Magna-Tiles",
                "kb://magna-tiles/Logiwa_IO_API_Carrier_Shipping_Option_Guide.docx",
            ),
            (
                MAGNA / "PO Receipt Implementation.pptx",
                "Magna-Tiles",
                "kb://magna-tiles/PO Receipt Implementation.pptx",
            ),
            (
                MAGNA / "Shipment_Inventory_Webhook_Guide.pdf",
                "Magna-Tiles",
                "kb://magna-tiles/Shipment_Inventory_Webhook_Guide.pdf",
            ),
        ]
    )

    articles = []
    seen = set()
    for path, origin, url in jobs:
        key = path.name.lower()
        if key in seen:
            continue
        seen.add(key)
        article = article_from(path, origin, url)
        if article:
            articles.append(article)

    OUT.write_text(json.dumps(articles, ensure_ascii=False, indent=2), encoding="utf-8")
    print(f"Wrote {len(articles)} articles to {OUT} ({OUT.stat().st_size} bytes)")


if __name__ == "__main__":
    main()
